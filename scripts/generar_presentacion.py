"""Genera la presentación PowerPoint de la investigación de mercado DuniPets."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import os

OUT_PATH = r"C:\Users\tuto6\Projects\DuniPets\reports\DuniPets_Investigacion_Mercado_2026.pptx"

# ── Paleta de color ──────────────────────────────────────────────────────────
VERDE        = RGBColor(0x2E, 0x7D, 0x32)   # verde oscuro — identidad
VERDE_CLARO  = RGBColor(0xA5, 0xD6, 0xA7)   # verde suave — fondos
GRIS_OSCURO  = RGBColor(0x21, 0x21, 0x21)   # casi negro — títulos
GRIS_MEDIO   = RGBColor(0x55, 0x55, 0x55)   # gris — cuerpo
BLANCO       = RGBColor(0xFF, 0xFF, 0xFF)
NARANJA      = RGBColor(0xE6, 0x5C, 0x00)   # acento — datos clave
FONDO_SLIDE  = RGBColor(0xF9, 0xFB, 0xF9)   # blanco roto

# ── Helpers ──────────────────────────────────────────────────────────────────
def set_bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_rect(slide, l, t, w, h, fill_color, line_color=None):
    shape = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
    else:
        shape.line.fill.background()
    return shape

def add_textbox(slide, text, l, t, w, h, size=12, bold=False, color=None,
                align=PP_ALIGN.LEFT, wrap=True):
    txb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    txb.word_wrap = wrap
    tf = txb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color or GRIS_OSCURO
    return txb

def add_label_value(slide, label, value, l, t, label_w=1.8, val_w=2.4,
                    row_h=0.32, label_size=9, val_size=10, val_color=None):
    add_textbox(slide, label, l, t, label_w, row_h, size=label_size,
                bold=True, color=GRIS_MEDIO)
    add_textbox(slide, value, l + label_w, t, val_w, row_h, size=val_size,
                color=val_color or GRIS_OSCURO)

def slide_header(slide, title, subtitle=None):
    """Banda superior verde con título blanco."""
    add_rect(slide, 0, 0, 10, 1.05, VERDE)
    add_textbox(slide, title, 0.35, 0.12, 9.3, 0.6,
                size=22, bold=True, color=BLANCO)
    if subtitle:
        add_textbox(slide, subtitle, 0.35, 0.68, 9.3, 0.32,
                    size=11, color=VERDE_CLARO)

def footnote(slide, text="DuniPets · Investigación de Mercado Bogotá · Junio 2026"):
    add_textbox(slide, text, 0.2, 7.1, 9.6, 0.28, size=7,
                color=GRIS_MEDIO, align=PP_ALIGN.CENTER)

def pill(slide, text, l, t, w=1.5, h=0.3, bg=VERDE_CLARO, fg=VERDE):
    r = add_rect(slide, l, t, w, h, bg)
    r.line.fill.background()
    tf = r.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = text
    run.font.size = Pt(9)
    run.font.bold = True
    run.font.color.rgb = fg
    # vertical centering
    from pptx.enum.text import MSO_ANCHOR
    tf.auto_size = None
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE

# ── Presentación ─────────────────────────────────────────────────────────────
prs = Presentation()
prs.slide_width  = Inches(10)
prs.slide_height = Inches(7.5)
blank = prs.slide_layouts[6]   # blank layout

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — PORTADA
# ════════════════════════════════════════════════════════════════════════════
s1 = prs.slides.add_slide(blank)
set_bg(s1, VERDE)

# Bloque central blanco
add_rect(s1, 1.0, 1.6, 8.0, 4.3, BLANCO)

# Logo/nombre grande
add_textbox(s1, "DuniPets", 1.3, 1.9, 7.4, 1.1,
            size=52, bold=True, color=VERDE, align=PP_ALIGN.CENTER)

# Línea decorativa
add_rect(s1, 3.5, 3.05, 3.0, 0.05, VERDE_CLARO)

# Subtítulo
add_textbox(s1, "Investigación de Mercado", 1.3, 3.2, 7.4, 0.55,
            size=20, bold=False, color=GRIS_OSCURO, align=PP_ALIGN.CENTER)
add_textbox(s1, "Marketplace Premium de Cuidado de Mascotas · Bogotá, Colombia",
            1.3, 3.75, 7.4, 0.45, size=12, color=GRIS_MEDIO, align=PP_ALIGN.CENTER)

# Píldoras de contexto
pill(s1, "Junio 2026",        2.5,  4.55, w=1.6)
pill(s1, "5 competidores",    4.2,  4.55, w=1.6)
pill(s1, "3 perfiles cliente", 5.9, 4.55, w=1.9)

# Tagline inferior
add_textbox(s1, "Cero jaulas · Entorno familiar · Paz mental para el dueño",
            1.3, 5.4, 7.4, 0.4, size=10, color=BLANCO,
            align=PP_ALIGN.CENTER)

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — MAPA COMPETITIVO
# ════════════════════════════════════════════════════════════════════════════
s2 = prs.slides.add_slide(blank)
set_bg(s2, FONDO_SLIDE)
slide_header(s2, "Mapa Competitivo", "Quién está activo · Quién salió del mercado · Junio 2026")

# ── Activos ──
add_textbox(s2, "COMPETIDORES ACTIVOS", 0.35, 1.2, 4.5, 0.3,
            size=9, bold=True, color=VERDE)
add_rect(s2, 0.35, 1.5, 4.5, 0.02, VERDE_CLARO)

activos = [
    ("Dog Garden",   "Guardería física · 3 sedes",   "Sin jaulas · Solo perros",  "Desde $35.000/día"),
    ("Mywak",        "App / Marketplace",             "GPS + cuidadores certificados Cruz Roja", "No publicado"),
    ("Diverpool",    "Centro integral premium",       "Hotel + spa + vet + gatos en sede", "$40K–$80K/noche"),
    ("Cruz Roja Bog.","Guardería / Hotel físico",     "20.000 m² · 15 años experiencia", "Desde $6.900/hora"),
    ("Tierra de Perros","Hotel campestre",            "Vía Choachi · Membresías",  "No publicado"),
]

for i, (nombre, tipo, diferenciador, precio) in enumerate(activos):
    y = 1.62 + i * 0.88
    add_rect(s2, 0.35, y, 4.5, 0.78, BLANCO, line_color=VERDE_CLARO)
    pill(s2, "ACTIVO", 0.42, y + 0.06, w=0.72, h=0.22,
         bg=RGBColor(0xE8, 0xF5, 0xE9), fg=VERDE)
    add_textbox(s2, nombre, 1.22, y + 0.04, 2.4, 0.28, size=11, bold=True, color=GRIS_OSCURO)
    add_textbox(s2, tipo,   1.22, y + 0.30, 2.4, 0.22, size=8,  color=GRIS_MEDIO)
    add_textbox(s2, diferenciador, 3.65, y + 0.04, 1.15, 0.36, size=7.5, color=GRIS_MEDIO)
    add_textbox(s2, precio, 3.65, y + 0.46, 1.15, 0.22, size=9, bold=True, color=VERDE)

# ── Inactivos ──
add_textbox(s2, "SALIERON DEL MERCADO", 5.15, 1.2, 4.5, 0.3,
            size=9, bold=True, color=NARANJA)
add_rect(s2, 5.15, 1.5, 4.5, 0.02, RGBColor(0xFF, 0xCC, 0xBC))

inactivos = [
    ("Waggy",          "App / Marketplace",  "waggy.co redirige a atom.com · Abandonado"),
    ("DogHero Colombia","App / Marketplace", "Sin presencia web activa en Colombia desde 2020"),
]

for i, (nombre, tipo, razon) in enumerate(inactivos):
    y = 1.62 + i * 1.0
    add_rect(s2, 5.15, y, 4.5, 0.82, BLANCO,
             line_color=RGBColor(0xFF, 0xCC, 0xBC))
    pill(s2, "INACTIVO", 5.22, y + 0.08, w=0.88, h=0.22,
         bg=RGBColor(0xFF, 0xEB, 0xE6), fg=NARANJA)
    add_textbox(s2, nombre, 6.18, y + 0.05, 2.5, 0.28, size=11, bold=True, color=GRIS_OSCURO)
    add_textbox(s2, tipo,   6.18, y + 0.32, 2.5, 0.22, size=8,  color=GRIS_MEDIO)
    add_textbox(s2, razon,  5.22, y + 0.52, 4.3, 0.26, size=8,  color=NARANJA)

# Insight box
add_rect(s2, 5.15, 3.75, 4.5, 1.3, RGBColor(0xF1, 0xF8, 0xE9), line_color=VERDE_CLARO)
add_textbox(s2, "Insight clave", 5.3, 3.83, 4.2, 0.28, size=9, bold=True, color=VERDE)
add_textbox(s2,
    "Los dos marketplaces de referencia (Waggy y DogHero) "
    "ya no operan en Colombia. El segmento de plataformas digitales "
    "de cuidado de mascotas en Bogotá está disponible — "
    "Mywak lo ocupa parcialmente sin comunicar calidad diferenciada.",
    5.3, 4.12, 4.2, 0.88, size=8.5, color=GRIS_OSCURO)

footnote(s2)

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — BENCHMARK DE PRECIOS
# ════════════════════════════════════════════════════════════════════════════
s3 = prs.slides.add_slide(blank)
set_bg(s3, FONDO_SLIDE)
slide_header(s3, "Benchmark de Precios", "Comparativo por servicio · Todos los precios en COP · Fuente: auditoría web junio 2026")

headers = ["Servicio", "Dog Garden", "Diverpool", "Cruz Roja", "Mywak", "DuniPets"]
col_x   = [0.25,        2.05,         3.45,         4.85,        6.25,    7.65]
col_w   = [1.7,         1.3,          1.3,          1.3,         1.3,     1.75]

rows = [
    ("Hospedaje\nnocturno",     "Por consultar",     "$40K–$80K/n",   "Por plan",   "No publicado",  "$47.500/noche"),
    ("Guardería\ndiurna",       "Desde $35.000",      "Por consultar", "$55.200 ~8h","No publicado",  "$35.000/día"),
    ("Paseo 1 hora",            "Incluido*",          "Por consultar", "No aplica",  "No publicado",  "$13.500"),
    ("Visita gatos\na domicilio","No ofrece",         "No ofrece",     "No ofrece",  "Sin confirmar", "$12.000–$18.000"),
    ("House sitting",           "No ofrece",          "No ofrece",     "No ofrece",  "Sin confirmar", "Por definir**"),
]

# Cabecera de tabla
add_rect(s3, 0.25, 1.15, 9.5, 0.38, VERDE)
for j, (hdr, cx, cw) in enumerate(zip(headers, col_x, col_w)):
    align = PP_ALIGN.CENTER if j > 0 else PP_ALIGN.LEFT
    add_textbox(s3, hdr, cx + 0.05, 1.18, cw - 0.1, 0.32,
                size=9.5, bold=True, color=BLANCO, align=align)

# Filas
for i, (servicio, *vals) in enumerate(rows):
    y = 1.55 + i * 1.04
    row_bg = BLANCO if i % 2 == 0 else RGBColor(0xF4, 0xF9, 0xF4)
    add_rect(s3, 0.25, y, 9.5, 0.98, row_bg, line_color=VERDE_CLARO)

    # Nombre del servicio
    add_textbox(s3, servicio, col_x[0] + 0.06, y + 0.08,
                col_w[0] - 0.1, 0.82, size=9.5, bold=True, color=GRIS_OSCURO)

    for j, (val, cx, cw) in enumerate(zip(vals, col_x[1:], col_w[1:])):
        is_dunipets = (j == 4)
        is_empty    = val in ("No ofrece", "No aplica", "Sin confirmar",
                               "No publicado", "Por consultar", "Por definir**")
        color = VERDE if is_dunipets else (RGBColor(0xBD, 0xBD, 0xBD) if is_empty else GRIS_OSCURO)
        bold  = is_dunipets
        size  = 10 if is_dunipets else 9
        add_textbox(s3, val, cx + 0.06, y + 0.22, cw - 0.1, 0.55,
                    size=size, bold=bold, color=color, align=PP_ALIGN.CENTER)

# Nota de fuente
add_textbox(s3,
    "* En Dog Garden los paseos están incluidos dentro del servicio de guardería, no se venden por separado.  "
    "** House sitting DuniPets: rango recomendado $65.000–$85.000/noche basado en referencias históricas.",
    0.25, 6.8, 9.5, 0.32, size=7, color=GRIS_MEDIO)

footnote(s3)

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — LOS 3 GAPS DE MERCADO
# ════════════════════════════════════════════════════════════════════════════
s4 = prs.slides.add_slide(blank)
set_bg(s4, FONDO_SLIDE)
slide_header(s4, "Los 3 Gaps de Mercado", "Oportunidades confirmadas por la auditoría de competidores · Junio 2026")

gaps = [
    (
        "01",
        "El mercado felino\nestá sin atender",
        "0 de 5 competidores activos ofrece visitas a domicilio "
        "para gatos. El gato necesita quedarse en su entorno; "
        "el mercado lo obliga a salir o a quedarse solo.",
        "DuniPets es el único jugador activo con visitas a domicilio "
        "para gatos en Bogotá. Precio de entrada $12.000 reduce barrera de primera prueba.",
    ),
    (
        "02",
        "\"Sin jaulas\" solo\nen instalaciones",
        "Dog Garden validó demanda con 3 sedes activas. "
        "Pero su modelo requiere traslado a sede física — "
        "solo cubre el norte de Bogotá.",
        "\"Hogar de verdad\" es el paso siguiente: el cuidador recibe "
        "en su casa o se queda en la del cliente. Sin sedes fijas, "
        "cobertura distribuida por toda la ciudad.",
    ),
    (
        "03",
        "Precios opacos\nen toda la industria",
        "4 de 5 competidores activos no publican precios. "
        "La norma es \"contáctenos por WhatsApp\". "
        "Genera abandono y desconfianza.",
        "DuniPets publica tabla completa de precios desde el primer "
        "contacto. Transparencia como ventaja de conversión — "
        "el referente de precio del mercado.",
    ),
]

for i, (num, titulo, evidencia, oportunidad) in enumerate(gaps):
    x = 0.22 + i * 3.25
    # Tarjeta
    add_rect(s4, x, 1.18, 3.1, 5.55, BLANCO, line_color=VERDE_CLARO)
    # Número grande
    add_rect(s4, x, 1.18, 3.1, 0.6, VERDE)
    add_textbox(s4, num, x + 0.08, 1.2, 0.55, 0.55,
                size=26, bold=True, color=BLANCO)

    add_textbox(s4, titulo, x + 0.12, 1.88, 2.88, 0.75,
                size=12.5, bold=True, color=GRIS_OSCURO)

    add_textbox(s4, "EVIDENCIA", x + 0.12, 2.72, 2.88, 0.22,
                size=7.5, bold=True, color=NARANJA)
    add_textbox(s4, evidencia, x + 0.12, 2.94, 2.88, 1.5,
                size=8.5, color=GRIS_OSCURO)

    add_rect(s4, x + 0.12, 4.5, 2.86, 0.02, VERDE_CLARO)

    add_textbox(s4, "CÓMO LO CAPITALIZA DUNIPETS", x + 0.12, 4.58, 2.88, 0.22,
                size=7.5, bold=True, color=VERDE)
    add_textbox(s4, oportunidad, x + 0.12, 4.8, 2.88, 1.5,
                size=8.5, color=GRIS_OSCURO)

footnote(s4)

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — PERFILES DE CLIENTE
# ════════════════════════════════════════════════════════════════════════════
s5 = prs.slides.add_slide(blank)
set_bg(s5, FONDO_SLIDE)
slide_header(s5, "Los 3 Perfiles de Cliente", "Arquetipos prioritarios para Bogotá · Junio 2026")

perfiles = [
    {
        "nombre": "Valentina",
        "arquetipo": "La Viajera\ncon Gatos",
        "zona": "Chapinero / Teusaquillo",
        "mascota": "1–2 gatos",
        "servicio": "Visitas domicilio\n+ House sitting",
        "barrera": "No sabe que el\nservicio existe",
        "canal": "Instagram · Grupos\nFacebook de gatos",
        "mensaje": "\"Tus gatos se quedan\nen su hogar.\nNosotros vamos\ndonde ellos.\"",
        "valor": "$700K–$1.05M / año",
        "gap": "Gap 1",
    },
    {
        "nombre": "Andrés / Camila",
        "arquetipo": "El Ejecutivo\nViajero",
        "zona": "Usaquén / Chicó",
        "mascota": "Perro mediano-grande",
        "servicio": "Hospedaje nocturno",
        "barrera": "Mala experiencia\nprevia en kennel",
        "canal": "Recomendación\ndirecta · Google",
        "mensaje": "\"Tu perro en un\nhogar de verdad.\nFotos cada noche,\ngarantizado.\"",
        "valor": "$1.52M / año",
        "gap": "Gap 2",
    },
    {
        "nombre": "Daniela / Juan Felipe",
        "arquetipo": "La Profesional\nAgenda Llena",
        "zona": "Suba / Cedritos / Chapinero",
        "mascota": "Perro pequeño-mediano",
        "servicio": "Guardería diurna\n+ Paseos",
        "barrera": "Fricción del proceso\nde cotización",
        "canal": "Instagram · TikTok\n· WhatsApp edificio",
        "mensaje": "\"$13.500 el paseo.\nRecogemos en\ntu puerta. Reserva\nen 2 minutos.\"",
        "valor": "$1.3M–$3.4M / año",
        "gap": "Gap 3",
    },
]

LABELS = ["ZONA", "MASCOTA", "SERVICIO PRINCIPAL", "BARRERA", "CANAL", "MENSAJE CLAVE", "VALOR ANUAL EST."]

for i, p in enumerate(perfiles):
    x = 0.22 + i * 3.25
    add_rect(s5, x, 1.18, 3.1, 5.6, BLANCO, line_color=VERDE_CLARO)

    # Cabecera de tarjeta
    add_rect(s5, x, 1.18, 3.1, 0.88, VERDE)
    pill(s5, p["gap"], x + 0.12, 1.24, w=0.72, h=0.22,
         bg=VERDE_CLARO, fg=VERDE)
    add_textbox(s5, p["nombre"], x + 0.9, 1.22, 2.1, 0.28,
                size=10, bold=True, color=BLANCO)
    add_textbox(s5, p["arquetipo"], x + 0.12, 1.5, 2.9, 0.5,
                size=12, bold=True, color=BLANCO)

    campos = [
        p["zona"], p["mascota"], p["servicio"],
        p["barrera"], p["canal"], p["mensaje"], p["valor"],
    ]
    for j, (label, valor) in enumerate(zip(LABELS, campos)):
        y = 2.14 + j * 0.73
        is_mensaje = label == "MENSAJE CLAVE"
        is_valor   = label == "VALOR ANUAL EST."
        add_textbox(s5, label, x + 0.12, y, 2.88, 0.2,
                    size=7, bold=True,
                    color=VERDE if not is_valor else NARANJA)
        add_textbox(s5, valor, x + 0.12, y + 0.2, 2.88, 0.48,
                    size=8.5 if not is_mensaje else 8,
                    color=GRIS_OSCURO if not is_valor else NARANJA,
                    bold=is_valor)

footnote(s5)

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — POSICIONAMIENTO Y PRÓXIMOS PASOS
# ════════════════════════════════════════════════════════════════════════════
s6 = prs.slides.add_slide(blank)
set_bg(s6, FONDO_SLIDE)
slide_header(s6, "Posicionamiento y Próximos Pasos", "Recomendaciones basadas en la investigación · Junio 2026")

# ── Posicionamiento ──
add_textbox(s6, "POSICIONAMIENTO RECOMENDADO", 0.35, 1.2, 5.8, 0.28,
            size=9, bold=True, color=VERDE)
add_rect(s6, 0.35, 1.5, 5.8, 0.03, VERDE_CLARO)

add_rect(s6, 0.35, 1.55, 5.8, 1.55, BLANCO, line_color=VERDE_CLARO)
add_textbox(s6,
    "DuniPets es el servicio de cuidado de mascotas más confiable "
    "de Bogotá para dueños que no aceptan jaulas. Cuidamos perros "
    "y gatos en entornos domésticos — en el hogar de nuestros "
    "cuidadores o en el tuyo — con precios claros y cero sorpresas.",
    0.5, 1.62, 5.5, 1.4, size=11, color=GRIS_OSCURO)

# Tres pilares
add_textbox(s6, "LOS TRES PILARES", 0.35, 3.22, 5.8, 0.28,
            size=9, bold=True, color=VERDE)
pilares = [
    ("Cero jaulas",          "Siempre en un hogar real"),
    ("Solo para mascotas",   "Perros y gatos · Juntos o separados"),
    ("Precio transparente",  "Publicado · Sin sorpresas · Sin WhatsApp previo"),
]
for i, (p_titulo, p_sub) in enumerate(pilares):
    y = 3.55 + i * 0.65
    add_rect(s6, 0.35, y, 5.8, 0.56, RGBColor(0xF1, 0xF8, 0xE9), line_color=VERDE_CLARO)
    add_rect(s6, 0.35, y, 0.08, 0.56, VERDE)
    add_textbox(s6, p_titulo, 0.55, y + 0.04, 2.5, 0.28, size=10, bold=True, color=VERDE)
    add_textbox(s6, p_sub,    0.55, y + 0.3,  5.4, 0.22, size=8.5, color=GRIS_MEDIO)

# ── Próximos pasos ──
add_textbox(s6, "PRÓXIMOS PASOS PRIORITARIOS", 6.45, 1.2, 3.3, 0.28,
            size=9, bold=True, color=NARANJA)
add_rect(s6, 6.45, 1.5, 3.3, 0.03, RGBColor(0xFF, 0xCC, 0xBC))

pasos = [
    ("S1–2", "Entrevistar 15 dueños de gatos",      "Validar Gap 1 con datos primarios"),
    ("S1",   "Mystery shopping Dog Garden",          "Mapear precios reales + experiencia"),
    ("S1",   "Descargar Mywak y simular reserva",   "Obtener precios y evaluar UX"),
    ("S1",   "Definir precio house sitting",         "Completar tabla de precios DuniPets"),
    ("S3–4", "Encuesta disposición a pagar",         "Validar sensibilidad por perfil"),
]
for i, (plazo, accion, objetivo) in enumerate(pasos):
    y = 1.58 + i * 1.06
    add_rect(s6, 6.45, y, 3.3, 0.92, BLANCO, line_color=RGBColor(0xFF, 0xCC, 0xBC))
    pill(s6, plazo, 6.52, y + 0.08, w=0.65, h=0.22,
         bg=RGBColor(0xFF, 0xEB, 0xE6), fg=NARANJA)
    add_textbox(s6, accion,   7.24, y + 0.06, 2.42, 0.28, size=9,   bold=True, color=GRIS_OSCURO)
    add_textbox(s6, objetivo, 6.52, y + 0.56, 3.14, 0.28, size=8,   color=GRIS_MEDIO)

footnote(s6)

# ── Guardar ──────────────────────────────────────────────────────────────────
os.makedirs(os.path.dirname(OUT_PATH), exist_ok=True)
prs.save(OUT_PATH)
print(f"Presentación guardada: {OUT_PATH}")
